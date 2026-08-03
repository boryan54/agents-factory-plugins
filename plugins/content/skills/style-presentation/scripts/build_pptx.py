# -*- coding: utf-8 -*-
"""Node-free сборка PPTX из картинок-слайдов (python-pptx). Для переносимого пакета.
Использование: python scripts/build_pptx.py <project.json>
Берёт out_dir из project.json, кладёт slide_01.png..slide_NN.png на слайды 16:9 во всю площадь."""
import sys, os, json, glob, re
from pptx import Presentation
from pptx.util import Inches

if len(sys.argv) < 2:
    sys.exit("usage: python build_pptx.py <project.json>")
cfg = json.load(open(sys.argv[1], encoding="utf-8-sig"))
out_dir = cfg["out_dir"]
name = cfg.get("project_name", "presentation")

pics = sorted(glob.glob(os.path.join(out_dir, "slide_*.png")),
              key=lambda p: int(re.search(r"slide_(\d+)", p).group(1)))
if not pics:
    sys.exit("нет картинок slide_*.png в " + out_dir)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for p in pics:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)

outpath = os.path.join(out_dir, name + ".pptx")
prs.save(outpath)
print("OK:", outpath, "(слайдов:", len(pics), ")")
