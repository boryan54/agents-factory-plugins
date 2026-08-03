# -*- coding: utf-8 -*-
"""
ШАБЛОН: редактируемый ТЕКСТОВЫЙ слайд PowerPoint (native, не картинка).

Когда использовать: пользователь просит слайд/презентацию в РЕДАКТИРУЕМОМ виде
(живой текст, не сгенерированная картинка), либо нужно 1:1 повторить оформление и
кегли существующего фирменного дека (таблицы, роадмапы, «Проблема-Решение» и т.п.).
Это ОТДЕЛЬНЫЙ режим от Genosai-генерации (см. SKILL.md, раздел «Режим B»).

Как снять стиль с образца (reference .pptx):
  1) Рендер слайдов в PNG для глаз (LibreOffice не нужен, есть PowerPoint COM):
     PowerShell:
       $pp=New-Object -ComObject PowerPoint.Application
       $pr=$pp.Presentations.Open($path,-1,0,0); $pr.Export($outDir,"PNG",1920,1080)
       $pr.Close(); $pp.Quit()
  2) Точные шрифты/кегли/цвета — python-pptx: пройти по slide.shapes -> text_frame
     -> paragraphs -> runs и прочитать run.font.name/size.pt/bold; для таблиц slide
     .shapes[i].table. ВАЖНО: узнай формат образца (prs.slide_width/height): деки часто
     идут в 2x (напр. 26.67x15") — тогда СТРОЙ новый дек В ТОМ ЖЕ размере, чтобы pt
     совпадали 1:1.

Грабли (уже учтены ниже):
  - Первый абзац авто-фигуры по умолчанию ЦЕНТРИРОВАН -> ставь PP_ALIGN.LEFT явно.
  - Висячий отступ пунктов -> pPr marL/indent (в EMU).
  - Высокая заполненная плашка полупустая -> vertical_anchor=MIDDLE (center=True).
  - Шрифт образца (напр. 'Eastman Medium') должен быть установлен, иначе подстановка.
  - Если .pptx открыт в PowerPoint/синкается Drive -> prs.save падает PermissionError:
    собирай во временную папку (env SLIDE_OUTDIR) или попроси закрыть файл.
  - Рендер QA и PDF: PowerPoint COM ($pres.Export ... "PNG"; $pres.SaveAs($pdf,32)).

Зависимости: pip install python-pptx ; PowerPoint (COM) для рендера/PDF.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- ПАЛИТРА/ШРИФТ образца (пример: набор rossiya-2050-eao) — подставь свои ---
TEAL = RGBColor(0x0F, 0x9D, 0xA6)   # акцент
DARK = RGBColor(0x1D, 0x3D, 0x42)   # основной текст
GREY_ROW = RGBColor(0xEE, 0xF1, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TXT = RGBColor(0x6B, 0x7B, 0x7E)
FONT = "Eastman Medium"
# Кегли образца ЕАО (дек 26.67x15"): H1=32, заголовок секции=22, пункты=19, таблица=20
SZ_TITLE, SZ_SUB, SZ_HDR, SZ_BODY, SZ_TABLE, SZ_CAP = 32, 18, 22, 19, 20, 14

SLIDE_W, SLIDE_H = 26.67, 15.0   # формат образца (16:9, 2x) — чтобы pt совпадали 1:1

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # пустой макет


def set_run(r, text, size, bold=False, color=DARK, font=FONT):
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font


def add_textbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    return tb, tf


def set_hanging(p, marL_in=0.32, indent_in=0.32):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(marL_in))))
    pPr.set('indent', str(-int(Inches(indent_in))))


def header_label(x, y, w, text):
    """teal-подзаголовок секции: точка + текст + короткая подчёркивающая линия."""
    tb, tf = add_textbox(x, y, w, 0.55)
    p = tf.paragraphs[0]
    set_run(p.add_run(), "●  ", 20, True, TEAL)
    set_run(p.add_run(), text, SZ_HDR, True, TEAL)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.03), Inches(y+0.6),
                                Inches(w-0.1), Inches(0.035))
    ln.fill.solid(); ln.fill.fore_color.rgb = TEAL; ln.line.fill.background()


def bullets_box(x, y, w, h, items, filled=False, gap=8, center=False):
    """items = [(жирный_лид, обычный_хвост), ...]. filled=True -> teal-плашка/белый текст."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    try: box.adjustments[0] = 0.03
    except Exception: pass
    if filled:
        box.fill.solid(); box.fill.fore_color.rgb = TEAL; box.line.fill.background()
        txt_c = dot_c = lead_c = WHITE
    else:
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = TEAL; box.line.width = Pt(1.25)
        txt_c = DARK; dot_c = lead_c = TEAL
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if center else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT                    # ГРАБЛИ: иначе 1-й абзац центр
        p.space_after = Pt(gap); p.space_before = Pt(2)
        set_hanging(p)                                 # висячий отступ
        set_run(p.add_run(), "●", 12, True, dot_c)
        set_run(p.add_run(), "   ", 12, False, dot_c)
        if lead: set_run(p.add_run(), lead, SZ_BODY, True, lead_c)
        if rest: set_run(p.add_run(), rest, SZ_BODY, False, txt_c)
    return box


def roadmap_table(x, y, w, h, rows):
    """rows = [(номер, описание, дата), ...]; teal-ячейка номера, чередование строк."""
    n = len(rows)
    tbl = slide.shapes.add_table(n, 3, Inches(x), Inches(y), Inches(w), Inches(h)).table
    tbl.first_row = False; tbl.horz_banding = False
    tbl.columns[0].width = Inches(0.7); tbl.columns[1].width = Inches(w-2.3)
    tbl.columns[2].width = Inches(1.6)
    for ri, (num, desc, date) in enumerate(rows):
        tbl.rows[ri].height = Inches(h/n)
        c0, c1, c2 = tbl.rows[ri].cells
        bg = GREY_ROW if ri % 2 == 0 else WHITE
        c0.fill.solid(); c0.fill.fore_color.rgb = TEAL
        c0.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = c0.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        set_run(pp.add_run(), num, SZ_HDR, True, WHITE)
        c1.fill.solid(); c1.fill.fore_color.rgb = bg
        c1.vertical_anchor = MSO_ANCHOR.MIDDLE
        c1.margin_left = Inches(0.18)
        set_run(c1.text_frame.paragraphs[0].add_run(), desc, SZ_TABLE, True, DARK)
        c2.fill.solid(); c2.fill.fore_color.rgb = bg
        c2.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = c2.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        set_run(pp.add_run(), date, SZ_TABLE, True, TEAL)


# ============ ПРИМЕР сборки (замени контент под задачу) ============
_tb, _tf = add_textbox(0.5, 0.15, SLIDE_W-0.9, 1.15)
set_run(_tf.paragraphs[0].add_run(), "ЗАГОЛОВОК СЛАЙДА", SZ_TITLE, True, TEAL)
set_run(_tf.add_paragraph().add_run(), "Подзаголовок слайда", SZ_SUB, False, GREY_TXT)

header_label(0.94, 1.12, 8.24, "Проблема")
bullets_box(0.57, 1.75, 8.24, 7.15, [
    ("Лид-ин: ", "текст пункта."),
    ("Ещё лид: ", "второй пункт."),
])
header_label(9.58, 1.12, 8.9, "Решение")
bullets_box(9.21, 1.75, 8.9, 12.63, [
    ("Пункт: ", "текст решения."),
], filled=True, gap=18, center=True)
header_label(18.87, 1.12, 7.66, "Дорожная карта")
roadmap_table(18.53, 1.8, 7.66, 12.4, [
    ("1", "Первый этап дорожной карты", "2026"),
    ("2", "Второй этап", "2027"),
])
_tb, _tf = add_textbox(0.57, SLIDE_H-0.5, 14, 0.4)
set_run(_tf.paragraphs[0].add_run(), "Источник: ...", SZ_CAP, False, GREY_TXT)

# Вывод: env SLIDE_OUTDIR позволяет собрать во временную папку, если целевой файл занят.
outdir = os.environ.get("SLIDE_OUTDIR", r"D:\GD\AiPrezy")
os.makedirs(outdir, exist_ok=True)
outpath = os.path.join(outdir, "text_slide.pptx")
prs.save(outpath)
print("saved:", outpath)
